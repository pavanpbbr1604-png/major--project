import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_ds_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette - Professional Data Science Teal & Slate
    NAVY = RGBColor(15, 23, 42)         # #0F172A
    SLATE = RGBColor(30, 41, 59)        # #1E293B
    TEAL = RGBColor(13, 148, 136)       # #0D9488
    CYAN = RGBColor(6, 182, 212)        # #06B6D4
    BLUE = RGBColor(37, 99, 235)        # #2563EB
    WHITE = RGBColor(255, 255, 255)     # #FFFFFF
    LIGHT_BG = RGBColor(248, 250, 252)  # #F8FAFC
    DARK_TEXT = RGBColor(15, 23, 42)

    def add_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="DATA SCIENCE PROJECT REVIEW"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEAL
        p.font.name = "Arial"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, accent_color=TEAL):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(226, 232, 240)
        shape.line.width = Pt(1)

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.font.name = "Arial"

        cb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), height - Inches(0.8))
        tf = cb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Arial"
            p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Data Science Theme)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s1, NAVY)

    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Automated Crowd Analytics & Density Estimation"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "A Data Science & Deep Learning Approach to Spatial Pattern Recognition"
    p2.font.size = Pt(20)
    p2.font.color.rgb = CYAN
    p2.font.name = "Arial"
    p2.space_before = Pt(12)

    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = SLATE
    card1.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(1.3), Inches(4.4), Inches(10.7), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "DATA SCIENCE PROJECT REVIEW"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.font.name = "Arial"

    p = tf.add_paragraph()
    p.text = "Presenter: Pavan BR | Team: Jeevan L, G Bharath Sai"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "Project Guide: Prof. Bibi Annie Oommen (Dept. of CSE)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.font.name = "Arial"
    p.space_before = Pt(4)

    p = tf.add_paragraph()
    p.text = "Department of Computer Science & Engineering | CMR Institute of Technology, Bengaluru"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.font.name = "Arial"
    p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 2: Problem & Data Science Objectives
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s2, LIGHT_BG)
    add_header(s2, "Data Science Problem Statement & Objectives")

    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "The Business Problem",
             ["Public safety management requires automated crowd monitoring to prevent stampedes.",
              "Manual counting cannot scale for real-time video surveillance feeds.",
              "Goal: Build an automated spatial data analytics system."])

    add_card(s2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Data Challenges",
             ["Extreme Scale Variance: Foreground vs background object resolution disparity.",
              "Spatial Overlap Noise: Inter-person occlusion corrupts standard bounding box metrics.",
              "Uncertainty: Blurry/dark photos cause false precision in hard scalar predictions."])

    add_card(s2, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "DS Project Deliverables",
             ["End-to-end ML data pipeline (Preprocessing ➔ Deep Inference ➔ Post-Processing).",
              "Subsampled spatial occupancy density grid mapping.",
              "Statistical uncertainty scoring & interactive analytics web application."])

    # -------------------------------------------------------------
    # SLIDE 3: Data Science Pipeline Architecture
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s3, LIGHT_BG)
    add_header(s3, "End-to-End Data Science Architecture")

    box_w, box_h = Inches(2.1), Inches(4.8)
    spacing = Inches(0.3)
    start_x = Inches(0.8)

    stages = [
        ("1. Data Ingestion", ["Image loading & normalization.", "Adaptive CLAHE contrast enhancement."]),
        ("2. Spatial Slicing", ["Grid-based tile slicing.", "Prevents downsampling resolution loss."]),
        ("3. Deep Learning", ["YOLOv8 feature extraction.", "Bounding box regression & confidence scoring."]),
        ("4. Custom Filtering", ["Containment promotion.", "IoM (Intersection over Min) fragment NMS."]),
        ("5. Density & UQ", ["Binary grid occupancy mapping.", "Statistical uncertainty quantification."])
    ]

    for idx, (stitle, sitems) in enumerate(stages):
        x_pos = start_x + idx * (box_w + spacing)
        add_card(s3, x_pos, Inches(1.8), box_w, box_h, stitle, sitems, accent_color=NAVY)

    # -------------------------------------------------------------
    # SLIDE 4: Data Preprocessing & Spatial Slicing
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s4, LIGHT_BG)
    add_header(s4, "Data Preprocessing & Scale-Invariant Slicing")

    add_card(s4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Adaptive Image Enhancement",
             ["Histogram Equalization: CLAHE (Contrast Limited Adaptive Histogram Equalization).",
              "Improves feature contrast in low-light, shadowed, or atmospheric blur conditions.",
              "Normalizes color channels for robust deep neural network feature extraction."])

    add_card(s4, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Grid Slicing (Tiled Inference)",
             ["Prevents Downsampling Loss: Resizing 4K frames to 640x640 destroys small distant objects.",
              "Spatial Slicing: Slices image into overlapping $T \\times T$ tiles with stride $S = T - O$.",
              "Native Feature Resolution: Evaluates background regions at full resolution.",
              "Coordinate Offset Remapping: Re-projects local detections to global spatial coordinates."])

    # -------------------------------------------------------------
    # SLIDE 5: Deep Learning & Custom Spatial NMS
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s5, LIGHT_BG)
    add_header(s5, "Deep Feature Extraction & Custom Spatial Filtering")

    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "YOLOv8 Feature Extraction",
             ["Model Architecture: YOLOv8 Single-Shot Detector executed via ONNX Runtime.",
              "Output Vector: Bounding box coordinates $[x_1, y_1, x_2, y_2]$ and class confidence scores $P(\\text{person})$.",
              "Fast execution optimized for real-time video surveillance feeds."])

    add_card(s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Custom Occlusion Post-Processing (IoM)",
             ["Pass 0 (Containment Promotion): Suppresses partial body sub-fragments (legs/torso) contained within full body candidate ($>60\\%$).",
              "Pass 1 (Intersection over Min IoM): Normalizes overlap by smaller box area alone: $\\text{IoM} = \\frac{\\text{Area}(b_1 \cap b_2)}{\\min(A_1, A_2)}$.",
              "Preserves adjacent individuals standing close together without false suppression."])

    # -------------------------------------------------------------
    # SLIDE 6: Spatial Density Mapping
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s6, LIGHT_BG)
    add_header(s6, "Subsampled Binary Grid Density Mapping")

    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "The Overcounting Problem",
             ["Summing bounding box areas overcounts people because overlapping boxes share pixels.",
              "Density map regression lacks discrete individual bounding boxes.",
              "Need an exact non-redundant spatial coverage metric."])

    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Subsampled Binary Grid Solution",
             ["Grid Downsampling ($s=8$): Discretizes bounding boxes onto a binary spatial matrix $M \\in \\{0,1\\}^{\\frac{H}{s} \\times \\frac{W}{s}}$.",
              "Bitwise-OR Union: Merges overlapping bounding box regions at the spatial grid cell level.",
              "Non-Redundant Area: $\\mathcal{A}_{\\text{occupied}} = s^2 \\cdot \\sum M[u,v]$.",
              "Density Ratio Classification: Categorizes crowd level into Low, Moderate, High, Extreme."])

    # -------------------------------------------------------------
    # SLIDE 7: Statistical Uncertainty Quantification
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s7, LIGHT_BG)
    add_header(s7, "Uncertainty Modeling & Dynamic Output Gating")

    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Multi-Factor Uncertainty Metric ($\mathcal{U}$)",
             ["Quantifies measurement risk using 4 diagnostic signals:",
              "  1. Mean Confidence Penalty ($1 - \\bar{C}$)",
              "  2. Small Object Disparity Ratio ($\\gamma_{\\text{small}}$)",
              "  3. Spatial Occlusion Density ($\\gamma_{\\text{occl}}$)",
              "  4. Tile Boundary Instability ($1 - \\mathcal{S}_{\\text{cons}}$)",
              "Composite Reliability Score: $\\mathcal{R}_{\\text{score}} = 1.0 - \\mathcal{U}$."])

    add_card(s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Statistical Output Gating",
             ["Prevents false precision in blurry, dark, or noisy photos.",
              "If strict diagnostic rules pass:",
              "  ➔ Formats output as Exact: 'Count = N'.",
              "If ANY diagnostic condition fails under noisy data:",
              "  ➔ Formats output as Lower-Bound Estimate: 'Count $\ge$ N'."])

    # -------------------------------------------------------------
    # SLIDE 8: Model Evaluation & Metrics
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s8, LIGHT_BG)
    add_header(s8, "Model Performance & Empirical Evaluation")

    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Ablation Experiment Analysis",
             ["Baseline YOLOv8 Model: MAE = 18.4 | Fragment Rate = 14.2% | Recall = 61.5%",
              "Baseline + Tiled Slicing: MAE = 7.2 | Fragment Rate = 11.8% | Recall = 91.0%",
              "Baseline + Containment NMS: MAE = 5.1 | Fragment Rate = 2.4% | Recall = 74.8%",
              "Full DS Proposed Pipeline: MAE = 1.6 | Fragment Rate = 1.1% | Recall = 94.8%"])

    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Latency & System Profiling",
             ["Pre-processing (CLAHE): 4.2 ms (3.0%)",
              "Tiled Deep Inference (4x 640x640): 118.5 ms (83.5%)",
              "Custom NMS & Containment: 11.2 ms (7.9%)",
              "Density Grid & Uncertainty Metric: 8.1 ms (5.6%)",
              "Total End-to-End Latency: 142.0 ms (~7 FPS real-time on CPU)."])

    # -------------------------------------------------------------
    # SLIDE 9: Web Application & Data Auditing
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s9, LIGHT_BG)
    add_header(s9, "Deployment Stack & Session Audit Database")

    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Software Stack & Architecture",
             ["Backend Framework: Python 3.12, Flask Web Server.",
              "ML Engine: ONNX Runtime CPU/GPU acceleration.",
              "Database: SQLite session history logging.",
              "Data Export Engine: ReportLab Executive PDF Generator & CSV Exporter."])

    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Interactive Web Dashboard",
             ["User interface for real-time image upload & processing.",
              "Interactive crowd density heatmaps & bounding box visualization.",
              "Historical session logging with session search & filtering.",
              "1-click executive PDF report & CSV dataset download."])

    # -------------------------------------------------------------
    # SLIDE 10: Conclusion & Data Science Impact
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s10, LIGHT_BG)
    add_header(s10, "Conclusion & Project Impact")

    add_card(s10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Key Project Achievements",
             ["Built an end-to-end Data Science & Deep Learning pipeline for crowd estimation.",
              "Fixed major object detection failure modes (scale variance & spatial occlusion noise).",
              "Developed non-redundant spatial grid density mapping & uncertainty scoring.",
              "Achieved significant MAE error reduction (18.4 ➔ 1.6) with real-time throughput."])

    add_card(s10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Future Scope in Data Science",
             ["Edge Deployment: Hardware acceleration on NVIDIA Jetson for 30+ FPS video streams.",
              "Time-Series Predictive Analytics: Forecasting crowd accumulation patterns over time.",
              "Automated Security Triggers: SMS / WhatsApp alert dispatch when density crosses safe limits."])

    # -------------------------------------------------------------
    # SLIDE 11: Thank You (Dark Theme)
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s11, NAVY)

    tb = s11.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Data Science Project Review & Discussion"
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEAL
    p2.font.name = "Arial"
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Automated Crowd Analytics & Density Estimation | CMRIT Bengaluru"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.font.name = "Arial"
    p3.space_before = Pt(20)
    p3.alignment = PP_ALIGN.CENTER

    out_path = os.path.join("docs", "Data_Science_Crowd_Analytics_Presentation.pptx")
    prs.save(out_path)
    print(f"Successfully saved Data Science presentation to: {out_path}")

if __name__ == "__main__":
    create_ds_presentation()
