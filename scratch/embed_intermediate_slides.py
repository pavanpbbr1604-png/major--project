import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

ppt_path = os.path.join("docs", "PPT", "BCS786 – Project Work_Review2.pptx")
prs = pptx.Presentation(ppt_path)

# Slide 12 is index 11 (Intermediate Output Screenshots)
s12 = prs.slides[11]

# Clear existing pictures or shapes added previously if needed
img1 = os.path.join("docs", "figures", "intermediate", "3_before_custom_nms.jpg")
img2 = os.path.join("docs", "figures", "intermediate", "4_after_custom_nms.jpg")
img3 = os.path.join("docs", "figures", "intermediate", "5_binary_occupancy_grid.jpg")

# Add pictures with captions
if os.path.exists(img1):
    s12.shapes.add_picture(img1, Inches(0.6), Inches(2.2), width=Inches(2.8))
if os.path.exists(img2):
    s12.shapes.add_picture(img2, Inches(3.6), Inches(2.2), width=Inches(2.8))
if os.path.exists(img3):
    s12.shapes.add_picture(img3, Inches(6.6), Inches(2.2), width=Inches(2.8))

# Update Slide 12 text content
if len(s12.shapes) >= 2 and s12.shapes[1].has_text_frame:
    tf = s12.shapes[1].text_frame
    tf.word_wrap = True
    tf.clear()
    
    p1 = tf.paragraphs[0]
    p1.text = "Data Flow: Raw Image → Preprocessing → Standard NMS (BEFORE) → Custom Containment NMS (AFTER) → Binary Grid Mask"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(30, 58, 138)
    
    p2 = tf.add_paragraph()
    p2.text = "• Left (BEFORE): Standard IoU NMS produces duplicate edge boxes and split torso/leg fragments."
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(51, 65, 85)
    
    p3 = tf.add_paragraph()
    p3.text = "• Middle (AFTER): Custom Containment Promotion & Fragment Merging joins split boxes into full-body detections."
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(51, 65, 85)
    
    p4 = tf.add_paragraph()
    p4.text = "• Right (GRID MASK): Subsampled 8x8 binary occupancy grid calculates exact ground density without double counting."
    p4.font.size = Pt(11)
    p4.font.color.rgb = RGBColor(51, 65, 85)

prs.save(ppt_path)
print(f"[INFO] Successfully updated Slide 12 in {ppt_path} with intermediate screenshots!")
