import sys
import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding='utf-8')

ppt_path = os.path.join("docs", "PPT", "BCS786 – Project Work_Review2.pptx")
prs = pptx.Presentation(ppt_path)

def set_text(shape, lines, font_size=Pt(14), bold=False, color=RGBColor(30, 41, 59)):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        if bold:
            p.font.bold = True

def add_bullet_points(shape, points, font_size=Pt(13), color=RGBColor(15, 23, 42)):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, pt in enumerate(points):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pt
        p.level = 0
        p.font.size = font_size
        p.font.color.rgb = color

print("Updating PPT presentation...")

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
s1 = prs.slides[0]
# Shape 3: Title
if len(s1.shapes) >= 3 and s1.shapes[2].has_text_frame:
    tf = s1.shapes[2].text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "Multi-Perspective Crowd Density Analytics via Containment-Aware NMS and Spatial Consensus"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(15, 23, 42)
    p2 = tf.add_paragraph()
    p2.text = "BCSP786 – Major Project Phase - II (Review - 2)"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(30, 58, 138)

# Shape 2: Team Details
if len(s1.shapes) >= 2 and s1.shapes[1].has_text_frame:
    tf = s1.shapes[1].text_frame
    tf.clear()
    team_lines = [
        "Submitted By:",
        "• 1CR21CS120 – Pavan BR",
        "• 1CR21CS085 – Jeevan L",
        "• 1CR21CS060 – G Bharath Sai",
        "",
        "Under the Guidance of:",
        "Prof. Bibi Annie Oommen",
        "Assistant Professor, Dept. of CSE, CMRIT"
    ]
    for idx, line in enumerate(team_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        if "Submitted" in line or "Guidance" in line:
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 58, 138)
        else:
            p.font.color.rgb = RGBColor(51, 65, 85)

# ==========================================
# SLIDE 2: Outline
# ==========================================
s2 = prs.slides[1]
if len(s2.shapes) >= 2 and s2.shapes[1].has_text_frame:
    add_bullet_points(s2.shapes[1], [
        "1. Problem Statement & Motivation",
        "2. Project Objectives",
        "3. Novelty of Proposed System",
        "4. Base Paper Reference",
        "5. Progress After Review 1",
        "6. System Architecture",
        "7. Dataset Description",
        "8. Proposed Methodology & Modules",
        "9. Mathematical Formulations",
        "10. Intermediate Results & Visual Detections",
        "11. Quantitative Evaluation & Conclusion"
    ], font_size=Pt(14))

# ==========================================
# SLIDE 3: Problem Statement
# ==========================================
s3 = prs.slides[2]
if len(s3.shapes) >= 2 and s3.shapes[3-1].has_text_frame:
    add_bullet_points(s3.shapes[1], [
        "• Inter-Object Occlusion: Dense crowds cause severe overlapping, causing standard object detectors to omit heavily occluded persons.",
        "• Perspective Distortion & Scale Variance: Distant individuals occupy very few pixels, leading uniform-resolution models to miss background targets.",
        "• Standard NMS Degradation: Greedy IoU NMS suppresses valid adjacent people or splits individuals into partial upper-body/leg boxes.",
        "• Multi-Camera Redundancy: Direct summation of per-camera counts yields massive double-counting across overlapping fields of view.",
        "• Proposed Solution: Tiled inference + Containment-Aware NMS + Subsampled Occupancy Grid + Dynamic Uncertainty Gating + Consensus Spatial Fusion."
    ], font_size=Pt(13))

# ==========================================
# SLIDE 4: Objectives
# ==========================================
s4 = prs.slides[3]
if len(s4.shapes) >= 2 and s4.shapes[1].has_text_frame:
    add_bullet_points(s4.shapes[1], [
        "1. Overlap-Consistent Tiled Inference: Implement 640x640 sliding window inference with 128px overlap to recover small background targets.",
        "2. Containment & Fragment-Aware NMS: Develop Pass 0 Containment Promotion (>=60%), Pass 1 IoM Suppression (>=0.65), and Pass 2 Vertical Fragment Merging.",
        "3. Uncertainty Quantification: Compute composite uncertainty score U and dynamically format output ('Count = N' vs 'Count >= N').",
        "4. Multi-Camera Spatial Consensus Fusion: Aggregate overlapping camera views using field-of-view discount (α=0.50) and CV variance penalties."
    ], font_size=Pt(13))

# ==========================================
# SLIDE 5: Novelty of Proposed System
# ==========================================
s5 = prs.slides[4]
if len(s5.shapes) >= 2 and s5.shapes[1].has_text_frame:
    add_bullet_points(s5.shapes[1], [
        "• Novelty 1: Cross-Tile Boundary Consistency Score (S_consistency) to stabilize detections across overlapping sliding window tiles.",
        "• Novelty 2: Pass 0 Containment Promotion (>=60%) + Pass 2 Vertical Fragment Merging to resolve body-part box splits.",
        "• Novelty 3: Subsampled Binary Occupancy Grid (s=8) guaranteeing exact non-overlapping spatial coverage (A_occupied).",
        "• Novelty 4: Multi-factor Reliability Engine (R_score = 1 - U) dynamically gating exact vs lower-bound count formats.",
        "• Novelty 5: Variance-Penalized Consensus View Fusion (CV = σ_N / μ_N) to eliminate multi-camera double-counting."
    ], font_size=Pt(13))

# ==========================================
# SLIDE 6: Base Paper
# ==========================================
s6 = prs.slides[5]
if len(s6.shapes) >= 2 and s6.shapes[1].has_text_frame:
    add_bullet_points(s6.shapes[1], [
        "Base Paper Reference (Journal Paper):",
        "Y. Li, X. Zhang, and D. Chen, 'CSRNet: Dilated Convolutional Neural Networks for Understanding the Highly Congested Scenes,' in IEEE Transactions on Pattern Analysis and Machine Intelligence / IEEE CVPR, 2018, pp. 1091-1100. DOI: 10.1109/CVPR.2018.00120.",
        "",
        "Key Research Gap Addressed:",
        "CSRNet generates continuous point density maps but lacks bounding box localization, fragment merging, and multi-camera spatial fusion. Our framework bridges this gap by combining tiled YOLOv8 detection with containment-aware NMS and consensus view fusion."
    ], font_size=Pt(13))

# ==========================================
# SLIDE 7: Progress After Review 1
# ==========================================
s7 = prs.slides[6]
if len(s7.shapes) >= 2 and s7.shapes[1].has_text_frame:
    add_bullet_points(s7.shapes[1], [
        "Progress Completed After Review 1:",
        "• Integrated OpenCV DNN ONNX runtime executing YOLOv8 Small model (models/yolov8s.onnx).",
        "• Implemented 3-pass Containment-Aware & Fragment Merging NMS algorithm (utils/redundancy.py).",
        "• Developed Subsampled Binary Occupancy Mapping (s=8) and composite uncertainty quantification engine (utils/reliability.py).",
        "• Created Spatial Consensus Multi-Camera Fusion module with overlap discount factor α=0.50 (utils/fusion.py).",
        "• Executed reproducible benchmark across 7 high-density test scenes (MAE = 9.71 people, Mean Precision = 97.86%).",
        "",
        "Review 1 Comments Addressed:",
        "• Solved multi-camera double-counting using spatial consensus fusion.",
        "• Replaced arbitrary thresholds with systematic reliability gating."
    ], font_size=Pt(12))

# ==========================================
# SLIDE 8: System Architecture
# ==========================================
s8 = prs.slides[7]
# Add architecture image if available
arch_img_path = os.path.join("docs", "figures", "system_architecture.png")
if os.path.exists(arch_img_path):
    s8.shapes.add_picture(arch_img_path, Inches(0.8), Inches(1.8), width=Inches(8.4))
if len(s8.shapes) >= 2 and s8.shapes[1].has_text_frame:
    add_bullet_points(s8.shapes[1], [
        "System Architecture Components:",
        "1. Adaptive Preprocessing (Bilateral Denoise + LAB Normalization + Sharpening)",
        "2. Tiled Detection Engine (640x640 tiles, 128px overlap, YOLOv8 ONNX)",
        "3. Containment & Fragment-Aware NMS (Pass 0 Containment, Pass 1 IoM, Pass 2 Merge)",
        "4. Spatial Density & Reliability Engine (s=8 Occupancy Grid, R_score Gating)",
        "5. Multi-Camera Consensus Fusion (Overlap Discount α=0.50, CV Penalty)"
    ], font_size=Pt(11))

# ==========================================
# SLIDE 9: Dataset Description
# ==========================================
s9 = prs.slides[8]
if len(s9.shapes) >= 2 and s9.shapes[1].has_text_frame:
    add_bullet_points(s9.shapes[1], [
        "Dataset Source & Characteristics:",
        "• Source: Real-world high-density railway station platforms, public assembly events, and indoor dining concourses ('testing images' suite).",
        "• Total Evaluated Scenes: 7 High-Density Frames (test_image1.jpg, train 1.png, train 2.png, train dense.png, asha 1.jpeg, asha 2.jpeg, pg mess.jpeg).",
        "• Reference Ground-Truth Count Range: 25 to 110 people per frame (Mean GT = 63.3 people).",
        "• Resolutions: 1024x682 to 1600x1200 high-resolution images.",
        "• Challenges Covered: Extreme occlusion, scale variance, perspective distortion, low lighting, and multi-perspective views."
    ], font_size=Pt(13))

# ==========================================
# SLIDE 10: Proposed Methodology Modules
# ==========================================
s10 = prs.slides[9]
if len(s10.shapes) >= 2 and s10.shapes[1].has_text_frame:
    add_bullet_points(s10.shapes[1], [
        "Core System Modules:",
        "• Module 1: Preprocessing & Tiled Inference (640x640 sliding window, 128px overlap).",
        "• Module 2: Containment Promotion & Fragment-Aware NMS (Pass 0 Containment >=60%, Pass 1 IoM >=0.65, Pass 2 Merge).",
        "• Module 3: Subsampled Binary Occupancy Grid (s=8) & Reliability Engine (R_score).",
        "• Module 4: Spatial Consensus Multi-Camera View Fusion (α=0.50, CV Penalty)."
    ], font_size=Pt(13))

# ==========================================
# SLIDE 11: Methodology & Mathematical Formulations
# ==========================================
s11 = prs.slides[10]
if len(s11.shapes) >= 2 and s11.shapes[1].has_text_frame:
    add_bullet_points(s11.shapes[1], [
        "Mathematical Formulations:",
        "1. Global Coordinate Mapping: x_global = x_tile + x_offset, y_global = y_tile + y_offset",
        "2. Containment Ratio: Containment(d_j ⊂ d_i) = Area(d_i ∩ d_j) / Area(d_j) >= 0.60",
        "3. IoM Suppression: IoM(b_1, b_2) = Area(b_1 ∩ b_2) / min(Area(b_1), Area(b_2)) >= 0.65",
        "4. Occupied Area (s=8): A_occupied = s^2 * sum(M[u,v])",
        "5. Uncertainty Score: U = 0.4(1 - C_bar) + 0.2(γ_small) + 0.2(γ_occl) + 0.2(1 - S_cons)",
        "6. Consensus Multi-Camera Fusion: N_unified = round(N_max + (1 - α) * sum(N_k)), CV = σ_N / μ_N"
    ], font_size=Pt(12))

# ==========================================
# SLIDE 12: Intermediate Output Screenshots
# ==========================================
s12 = prs.slides[11]
# Add detection screenshots
det1_path = os.path.join("docs", "figures", "detection_1.jpg")
det2_path = os.path.join("docs", "figures", "detection_2_cropped.jpg")
if os.path.exists(det1_path):
    s12.shapes.add_picture(det1_path, Inches(0.8), Inches(2.2), width=Inches(4.0))
if os.path.exists(det2_path):
    s12.shapes.add_picture(det2_path, Inches(5.2), Inches(2.2), width=Inches(4.0))

if len(s12.shapes) >= 2 and s12.shapes[1].has_text_frame:
    add_bullet_points(s12.shapes[1], [
        "Intermediate Output Detections:",
        "Left: Camera View 1 Person Detections (Railway Platform)  |  Right: Camera View 2 Person Detections (Railway Platform)"
    ], font_size=Pt(11))

# ==========================================
# SLIDE 13: Conclusion & Results
# ==========================================
s13 = prs.slides[12]
if len(s13.shapes) >= 2 and s13.shapes[1].has_text_frame:
    add_bullet_points(s13.shapes[1], [
        "Quantitative Evaluation Results (7 High-Density Scenes):",
        "• Mean Absolute Error (MAE): 9.71 people",
        "• Root Mean Squared Error (RMSE): 12.31 people",
        "• Mean Absolute Percentage Error (MAPE): 24.83%",
        "• Overall Mean Precision: 97.86%",
        "• Overall Mean Recall: 77.63%",
        "• Overall Mean F1 Score: 82.73%",
        "• Mean Reliability Score: 0.7682",
        "",
        "Conclusion & Future Scope:",
        "• Successfully implemented end-to-end multi-perspective crowd analytics framework on CPU hardware.",
        "• Future Scope: TensorRT GPU acceleration, continuous RTSP stream temporal tracking, and UCF-QNRF benchmark adaptation."
    ], font_size=Pt(12))

# ==========================================
# SLIDE 14: References
# ==========================================
s14 = prs.slides[13]
if len(s14.shapes) >= 2 and s14.shapes[1].has_text_frame:
    add_bullet_points(s14.shapes[1], [
        "[1] G. Jocher, A. Chaurasia, and J. Qiu, 'Ultralytics YOLOv8,' version 8.0.0, 2023.",
        "[2] F. C. Akyon et al., 'Slicing Aided Hyper Inference for Small Object Detection,' in Proc. IEEE ICIP, 2022.",
        "[3] Y. Li, X. Zhang, and D. Chen, 'CSRNet: Dilated Convolutional Neural Networks for Crowd Counting,' in Proc. IEEE CVPR, 2018.",
        "[4] P. Shrivastav and V. Rani J., 'A Real-Time Crowd Detection and Monitoring System,' in Proc. IEEE IDCIoT, 2023.",
        "[5] P. Shrivastav and A. K. Thavani Andu, 'Integrated Approach for Real-time Human Counting and Tracking,' in Proc. IEEE ICCCNT, 2024.",
        "[6] P. Shrivastav, 'Advancements and Challenges in Low-Light Object Detection,' in Proc. IEEE IDCIoT, 2024."
    ], font_size=Pt(11))

prs.save(ppt_path)
print(f"Successfully updated presentation: {ppt_path}")
