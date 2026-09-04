import sys
import os
import pptx

sys.stdout.reconfigure(encoding='utf-8')
ppt_path = os.path.join("docs", "PPT", "BCS786 – Project Work_Review2.pptx")
prs = pptx.Presentation(ppt_path)

print(f"Total Slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    title_text = slide.shapes.title.text if slide.shapes.title else "No Title"
    print(f"\n==========================================")
    print(f"SLIDE {i+1}: {title_text}")
    print(f"==========================================")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"  Shape {j+1}:")
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    print(f"    - {paragraph.text.strip()}")
