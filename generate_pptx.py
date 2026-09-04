import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    NAVY = RGBColor(15, 23, 42)       # #0F172A
    SLATE = RGBColor(30, 41, 59)      # #1E293B
    CYAN = RGBColor(6, 182, 212)      # #06B6D4
    BLUE = RGBColor(37, 99, 235)      # #2563EB
    WHITE = RGBColor(255, 255, 255)   # #FFFFFF
    LIGHT_BG = RGBColor(248, 250, 252) # #F8FAFC
    DARK_TEXT = RGBColor(15, 23, 42)
    GRAY_TEXT = RGBColor(71, 85, 105)

    def add_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="MAJOR PROJECT PRESENTATION"):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BLUE
        p.font.name = "Arial"

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, accent_color=BLUE):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(226, 232, 240)
        shape.line.width = Pt(1)

        # Title
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.font.name = "Arial"

        # Content
        cb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), height - Inches(0.8))
        tf = cb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Arial"
            p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s1, NAVY)

    # Title box
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multi-Perspective Crowd Density Analytics"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Occlusion-Aware Detection, Reliability Quantification & Spatial Consensus Fusion"
    p2.font.size = Pt(20)
    p2.font.color.rgb = CYAN
    p2.font.name = "Arial"
    p2.space_before = Pt(12)

    # Team details card
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = SLATE
    card1.line.fill.background()

    tb = s1.shapes.add_textbox(Inches(1.3), Inches(4.4), Inches(10.7), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "PROJECT AUTHORS & INSTITUTION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = "Arial"

    p = tf.add_paragraph()
    p.text = "Presented by: Pavan BR | Jeevan L | G Bharath Sai"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "Project Guide: Prof. Bibi Annie Oommen (Dept. of CSE)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.font.name = "Arial"
    p.space_before = Pt(4)

    p = tf.add_paragraph()
    p.text = "Department of Computer Science & Engineering | CMR Institute of Technology, Bengaluru"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(148, 163, 184)
    p.font.name = "Arial"
    p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 2: Project Motivation & Need
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s2, LIGHT_BG)
    add_header(s2, "Motivation & Project Significance")

    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Public Safety Risk",
             ["High-density public venues (railway platforms, stadiums, festivals) face severe stampede hazards.",
              "Manual head counting by security personnel is too slow and error-prone during surges.",
              "Need for automated early-warning computer vision systems."])

    add_card(s2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Technical Gap",
             ["Standard AI models fail when people block each other (occlusion).",
              "Distant background people get missed due to image downscaling.",
              "Multiple CCTV cameras cause severe double-counting in shared fields of view."])

    add_card(s2, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Project Vision",
             ["Build an end-to-end Machine Learning pipeline using YOLOv8 ONNX acceleration.",
              "Integrate custom post-processing algorithms for occlusion, tiling, and spatial fusion.",
              "Deploy an operational web dashboard with real-time analytics & PDF export."])

    # -------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s3, LIGHT_BG)
    add_header(s3, "Four Core Limitations of Standard Object Detectors")

    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4),
             "1. Feature Resolution Loss (Scale Variance)",
             ["Resizing 4K image feeds to standard 640x640 network inputs causes distant background people to shrink below receptive field thresholds.",
              "Result: High false-negative dropouts in background crowd regions."])

    add_card(s3, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.4),
             "2. NMS Occlusion Dropout & Fragmentation",
             ["Standard IoU-based Non-Maximum Suppression deletes valid overlapping people.",
              "Fragmented body boxes (legs or torso) get counted as extra separate individuals."])

    add_card(s3, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4),
             "3. Misleading Hard Counts (False Precision)",
             ["Traditional AI outputs rigid hard numbers (e.g., '142 people') even in dark/blurry scenes.",
              "Lacks uncertainty measurement to signal when visual data is ambiguous."])

    add_card(s3, Inches(6.8), Inches(4.5), Inches(5.6), Inches(2.4),
             "4. Multi-Camera Redundancy & Double Counting",
             ["Summing counts from 2+ overlapping CCTV angles causes massive double-counting.",
              "Lacks cross-view spatial consensus to merge shared coverage areas."])

    # -------------------------------------------------------------
    # SLIDE 4: Proposed System Architecture
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s4, LIGHT_BG)
    add_header(s4, "End-to-End System Pipeline & Architecture")

    box_w, box_h = Inches(2.1), Inches(4.8)
    spacing = Inches(0.3)
    start_x = Inches(0.8)

    stages = [
        ("Step 1: Tiled Crop", ["Input frame split into overlapping grid tiles.", "Preserves high-resolution features."]),
        ("Step 2: YOLOv8 ONNX", ["Deep Learning feature extraction per tile.", "Detects bounding boxes & confidence scores."]),
        ("Step 3: Containment NMS", ["Pass 0: Full-body promotion.", "Pass 1: IoM fragment suppression."]),
        ("Step 4: Reliability Score", ["Evaluates blur, lighting & density.", "Outputs exact or lower-bound count."]),
        ("Step 5: Spatial Fusion", ["Combines multi-camera views.", "Applies variance penalty for overlap."])
    ]

    for idx, (stitle, sitems) in enumerate(stages):
        x_pos = start_x + idx * (box_w + spacing)
        add_card(s4, x_pos, Inches(1.8), box_w, box_h, stitle, sitems, accent_color=NAVY)

    # -------------------------------------------------------------
    # SLIDE 5: Innovation 1 - Tiled Inference
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s5, LIGHT_BG)
    add_header(s5, "Innovation 1: Overlap-Consistent Tiled Inference")

    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "How Tiled Slicing Works",
             ["Image $I \in \mathbb{R}^{H \\times W}$ is partitioned into overlapping tiles of size $T \\times T$.",
              "Stride $S = T - O$ ensures overlap region $O$ prevents body boundary clipping.",
              "Each tile is evaluated at native resolution, restoring small background target recall.",
              "Localized predictions are offset-mapped back to full image global coordinates."])

    add_card(s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Tile Consistency Score ($\mathcal{S}_{\text{consistency}}$)",
             ["Quantifies prediction stability across adjacent overlapping tile boundaries.",
              "Formula checks if detections in overlap region agree with IoU $\ge 0.40$.",
              "High consistency ($>0.80$) confirms robust boundary matching.",
              "Low consistency triggers uncertainty penalty in final reliability metric."])

    # -------------------------------------------------------------
    # SLIDE 6: Innovation 2 - Containment & Fragment-Aware NMS
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s6, LIGHT_BG)
    add_header(s6, "Innovation 2: Containment Promotion & Fragment-Aware NMS")

    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Pass 0: Full-Body Containment Promotion",
             ["Calculates Containment Ratio: $\\frac{\\text{Area}(d_i \cap d_j)}{A_j}$.",
              "Identifies sub-fragment boxes (e.g. legs/torso) contained inside a full body ($>60\%$).",
              "Promotes full-body candidate by inheriting peak confidence $\max(c_i, c_j)$.",
              "Suppresses inner body fragment prior to confidence sorting."])

    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Pass 1: IoM (Intersection-over-Min) Filtering",
             ["Standard IoU divides by total union area, failing when people stand side-by-side.",
              "IoM computes: $\\frac{\\text{Area}(b_1 \cap b_2)}{\\min(A_1, A_2)}$.",
              "Normalizes by smaller box area alone, suppressing partial sub-boxes ($\ge 0.65$).",
              "Preserves adjacent overlapping individuals without false suppression."])

    # -------------------------------------------------------------
    # SLIDE 7: Innovation 3 - Reliability & Uncertainty Scoring
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s7, LIGHT_BG)
    add_header(s7, "Innovation 3: Uncertainty Quantification & Dynamic Count Gating")

    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Multi-Factor Uncertainty Score ($\mathcal{U}$)",
             ["Combines 4 diagnostic risk signals:",
              "  1. Mean Detection Confidence ($1 - \\bar{C}$)",
              "  2. Small-Object Ratio ($\\gamma_{\\text{small}}$)",
              "  3. Spatial Occlusion Ratio ($\\gamma_{\\text{occl}}$)",
              "  4. Tile Boundary Instability ($1 - \\mathcal{S}_{\\text{cons}}$)",
              "Overall Reliability Score $\mathcal{R}_{\\text{score}} = 1.0 - \\mathcal{U}$."])

    add_card(s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Dynamic Count Output Gating",
             ["Prevents false precision in blurry, dark, or dense scenes.",
              "If strict rules pass ($\mathcal{U} \le 0.40$, $\\bar{C} > 0.65$, $\\mathcal{S}_{\\text{cons}} > 0.80$):",
              "  ➔ Formats output as Exact: 'Count = N'.",
              "If ANY condition fails under noisy environments:",
              "  ➔ Formats output as Conservative Lower-Bound: 'Count $\ge$ N'."])

    # -------------------------------------------------------------
    # SLIDE 8: Innovation 4 - Multi-Camera Spatial Consensus Fusion
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s8, LIGHT_BG)
    add_header(s8, "Innovation 4: Multi-Camera Perspective Consensus Fusion")

    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Cross-View Spatial Unification",
             ["Given $K$ camera feeds with individual counts $N_1, N_2, \dots, N_K$.",
              "Establishes baseline anchor count $N_{\\max} = \max_k N_k$.",
              "Applies overlap factor $\\alpha \in [0.0, 1.0]$ to subtract shared field of view:",
              "  $N_{\\text{unified}} = \\lfloor N_{\\max} + (1.0 - \\alpha) \sum_{k \\neq \\text{max}} N_k \\rfloor$."])

    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Discrepancy Variance Penalty",
             ["Computes Coefficient of Variation across camera counts: $\\text{CV} = \\frac{\\sigma_N}{\\mu_N}$.",
              "Applies variance penalty to downweight confidence if one camera reports anomalous count dropouts.",
              "Final Fusion Confidence: $\\mathcal{F}_{\\text{confidence}} = \\bar{\\mathcal{R}}_{\\text{views}} \cdot (1.0 - \\text{Penalty})$.",
              "Ensures stable venue-wide count estimation."])

    # -------------------------------------------------------------
    # SLIDE 9: Web Application & System Implementation
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s9, LIGHT_BG)
    add_header(s9, "Production Web Dashboard & Deployment Stack")

    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Technology Stack",
             ["Backend Framework: Python 3.12, Flask Web Server.",
              "Inference Engine: ONNX Runtime CPU/GPU execution provider.",
              "Deep Learning Model: YOLOv8s ONNX weights.",
              "Image Processing: OpenCV, NumPy, Supervision.",
              "Database Audit: SQLite history tracking.",
              "Report Generator: ReportLab Executive PDF Engine."])

    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "User Interface Features",
             ["Single & Multi-Perspective Image Upload.",
              "Real-time visual bounding box annotation.",
              "Subsampled binary grid crowd density heatmaps.",
              "Historical session audit logs with CSV export.",
              "Automated 1-click executive PDF report downloads."])

    # -------------------------------------------------------------
    # SLIDE 10: Experimental Evaluation & Benchmarks
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s10, LIGHT_BG)
    add_header(s10, "Empirical Evaluation & Ablation Benchmarks")

    add_card(s10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Ablation Study Results",
             ["Baseline YOLOv8: MAE = 18.4 | Fragment Rate = 14.2% | Recall = 61.5%",
              "Baseline + Tiling: MAE = 7.2 | Fragment Rate = 11.8% | Recall = 91.0%",
              "Baseline + Containment: MAE = 5.1 | Fragment Rate = 2.4% | Recall = 74.8%",
              "Full Proposed Pipeline: MAE = 1.6 | Fragment Rate = 1.1% | Recall = 94.8%"])

    add_card(s10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Latency & System Profiling",
             ["Adaptive CLAHE Preprocessing: 4.2 ms (3.0%)",
              "Tiled Inference (4x 640x640): 118.5 ms (83.5%)",
              "Containment & IoM NMS: 11.2 ms (7.9%)",
              "Binary Density & Reliability Engine: 8.1 ms (5.6%)",
              "Total End-to-End Latency: 142.0 ms (~7 FPS on CPU)."])

    # -------------------------------------------------------------
    # SLIDE 11: Conclusion & Future Scope
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s11, LIGHT_BG)
    add_header(s11, "Conclusion & Future Scope")

    add_card(s11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Key Achievements",
             ["Successfully solved scale-dependent dropouts and body part fragmentation in dense crowds.",
              "Introduced dynamic uncertainty gating to eliminate false precision in blurry/dark scenes.",
              "Implemented multi-camera consensus fusion to prevent cross-view double counting.",
              "Delivered a production-ready web dashboard with database logs and PDF reporting."])

    add_card(s11, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Future Work",
             ["Edge Hardware Acceleration: Deploying on NVIDIA Jetson for 30+ FPS video processing.",
              "Live RTSP Video Stream Integration: Connecting directly to CCTV camera feeds.",
              "Automated Security Alerts: Instant SMS / WhatsApp alert dispatch during crowd surges.",
              "3D Camera Calibration: Homography estimation for precise spatial localization."])

    # -------------------------------------------------------------
    # SLIDE 12: Thank You / Q&A (Dark Theme)
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s12, NAVY)

    tb = s12.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(24)
    p2.font.color.rgb = CYAN
    p2.font.name = "Arial"
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Multi-Perspective Crowd Density Analytics | CMR Institute of Technology, Bengaluru"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.font.name = "Arial"
    p3.space_before = Pt(20)
    p3.alignment = PP_ALIGN.CENTER

    out_path = os.path.join("docs", "Multi_Perspective_Crowd_Analytics_Presentation.pptx")
    prs.save(out_path)
    print(f"Successfully saved presentation to: {out_path}")

if __name__ == "__main__":
    create_presentation()
